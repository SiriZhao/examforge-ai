from collections.abc import Callable
from pathlib import Path

from PIL import Image

from app.schemas.review import OCRConfig, ParsedFile, ParsedPage
from app.services.ocr_cache import load_ocr_cache, save_ocr_cache
from app.services.ocr_service import OCRError, run_ocr_on_image, run_ocr_on_path
from app.services.runtime_paths import find_poppler_path
from app.services.subprocess_utils import hide_subprocess_windows

SUPPORTED_PARSE_EXTENSIONS = {".pptx", ".pdf", ".docx", ".md", ".txt", ".png", ".jpg", ".jpeg"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg"}
PDF_TEXT_MIN_CHARS = 20


class ParseError(RuntimeError):
    pass


ProgressCallback = Callable[[str, float | None], None]


def parse_file(
    path: Path,
    ocr_config: OCRConfig | None = None,
    progress_callback: ProgressCallback | None = None,
) -> ParsedFile:
    config = ocr_config or OCRConfig()
    suffix = path.suffix.lower()
    warnings: list[str] = []
    cache_used = False

    if suffix not in SUPPORTED_PARSE_EXTENSIONS:
        raise ParseError(f"不支持解析该文件类型：{suffix or '未知'}。")
    if not path.exists():
        raise ParseError(f"文件不存在：{path.name}。")

    try:
        if suffix == ".docx":
            pages = parse_docx(path)
        elif suffix == ".pptx":
            pages = parse_pptx(path)
        elif suffix == ".md":
            pages = parse_markdown(path)
        elif suffix == ".txt":
            pages = parse_text(path)
        elif suffix == ".pdf":
            pages, warnings, cache_used = parse_pdf(path, config, progress_callback)
        else:
            if progress_callback:
                progress_callback(f"正在识别图片：{path.name}", None)
            pages = parse_image(path, config)
            warnings = []
            cache_used = False
    except (ParseError, OCRError):
        raise
    except Exception as exc:
        raise ParseError(f"解析 {path.name} 失败：{exc}") from exc

    raw_text = "\n\n".join(page.text for page in pages if page.text.strip())
    return ParsedFile(
        filename=path.name,
        file_type=suffix,
        path=str(path.resolve()),
        pages=pages,
        raw_text=raw_text,
        warnings=warnings,
        ocr_cache_used=cache_used,
    )


def parse_docx(path: Path) -> list[ParsedPage]:
    try:
        from docx import Document
    except ImportError as exc:
        raise ParseError("未安装 python-docx。") from exc

    document = Document(path)
    lines = [paragraph.text.strip() for paragraph in document.paragraphs]
    text = "\n".join(line for line in lines if line)
    return [ParsedPage(page_number=1, text=text, source="text_extract")]


def parse_pptx(path: Path) -> list[ParsedPage]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ParseError("未安装 python-pptx。") from exc

    presentation = Presentation(path)
    pages: list[ParsedPage] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_text.append(text)
        pages.append(
            ParsedPage(
                page_number=index,
                text="\n".join(slide_text),
                source="text_extract",
            )
        )
    return pages


def parse_markdown(path: Path) -> list[ParsedPage]:
    text = path.read_text(encoding="utf-8")
    return [ParsedPage(page_number=1, text=text, source="text_extract")]


def parse_text(path: Path) -> list[ParsedPage]:
    try:
        text = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        text = path.read_text(encoding="utf-8", errors="ignore")
    return [ParsedPage(page_number=1, text=text, source="text_extract")]


def parse_pdf(
    path: Path,
    ocr_config: OCRConfig,
    progress_callback: ProgressCallback | None = None,
) -> tuple[list[ParsedPage], list[str], bool]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ParseError("未安装 pypdf。") from exc

    if progress_callback:
        progress_callback("正在分析文件类型：PDF", 0.02)

    cached_pages = load_ocr_cache(path, ocr_config)
    if cached_pages:
        if progress_callback:
            progress_callback("已使用 OCR 缓存。", 0.12)
        return cached_pages, ["已使用 OCR 缓存。"], True

    try:
        reader = PdfReader(str(path))
    except Exception as exc:
        raise ParseError(f"无法打开 PDF：{path.name}。") from exc

    pages: list[ParsedPage] = []
    warnings: list[str] = []
    total_pages = len(reader.pages)
    text_pages: dict[int, str] = {}
    ocr_page_numbers: list[int] = []

    if progress_callback:
        progress_callback("正在检查 PDF 文本层。", 0.04)

    for index, page in enumerate(reader.pages, start=1):
        if progress_callback:
            progress_callback(f"正在检查 PDF 文本层：第 {index}/{total_pages} 页", min(0.2, index / max(total_pages, 1) * 0.2))
        text = (page.extract_text() or "").strip()
        if len(text) >= PDF_TEXT_MIN_CHARS:
            text_pages[index] = text
        else:
            ocr_page_numbers.append(index)

    if not ocr_page_numbers:
        if progress_callback:
            progress_callback("已检测到文字版 PDF，跳过 OCR。", 0.35)
        return [
            ParsedPage(page_number=index, text=text_pages.get(index, ""), source="text_extract")
            for index in range(1, total_pages + 1)
        ], warnings, False

    selected_ocr_pages = select_ocr_pages(ocr_page_numbers, total_pages, ocr_config.mode)
    skipped = [page for page in ocr_page_numbers if page not in selected_ocr_pages]
    if skipped and progress_callback:
        progress_callback(f"快速模式已跳过 {len(skipped)} 个低信息扫描页，可切换完整模式全量 OCR。", 0.3)

    ocr_text_by_page: dict[int, str] = {}
    for ocr_index, page_number in enumerate(selected_ocr_pages, start=1):
        if progress_callback:
            progress_callback(
                f"正在 OCR 第 {ocr_index}/{len(selected_ocr_pages)} 页（PDF 第 {page_number} 页）。",
                0.3 + (ocr_index / max(len(selected_ocr_pages), 1)) * 0.45,
            )
        try:
            image = render_pdf_page_to_image(path, page_number)
            ocr_text_by_page[page_number] = run_ocr_on_image(image, ocr_config)
        except (ParseError, OCRError) as exc:
            warning = f"PDF 第 {page_number} 页 OCR 失败，已跳过该页：{exc}"
            warnings.append(warning)
            ocr_text_by_page[page_number] = ""

    if progress_callback:
        progress_callback("正在清洗 OCR 文本。", 0.78)

    for index in range(1, total_pages + 1):
        if index in text_pages:
            pages.append(ParsedPage(page_number=index, text=text_pages[index], source="text_extract"))
        elif index in ocr_text_by_page:
            warning = next((item for item in warnings if f"第 {index} 页" in item), None)
            pages.append(ParsedPage(page_number=index, text=ocr_text_by_page[index], source="ocr_fallback", warning=warning))
        else:
            pages.append(ParsedPage(page_number=index, text="", source="ocr_fallback", warning="快速模式跳过该扫描页。"))

    if selected_ocr_pages:
        save_ocr_cache(path, ocr_config, pages, page_count=total_pages)

    return pages, warnings, False


def select_ocr_pages(page_numbers: list[int], total_pages: int, mode: str) -> list[int]:
    if mode == "full" or len(page_numbers) <= 8:
        return page_numbers
    priority = set(page_numbers[:6])
    priority.update(page for page in page_numbers if page <= 3)
    priority.update(page for page in page_numbers if page % 5 == 0)
    priority.update(page for page in page_numbers if page >= max(total_pages - 1, 1))
    return [page for page in page_numbers if page in priority][:12]


def parse_image(path: Path, ocr_config: OCRConfig) -> list[ParsedPage]:
    text = run_ocr_on_path(path, ocr_config)
    return [ParsedPage(page_number=1, text=text, source="ocr_fallback")]


def render_pdf_page_to_image(path: Path, page_number: int) -> Image.Image:
    try:
        from pdf2image import convert_from_path
    except ImportError as exc:
        raise ParseError("未安装 pdf2image。") from exc

    try:
        poppler_path = find_poppler_path()
        with hide_subprocess_windows():
            images = convert_from_path(
                str(path),
                first_page=page_number,
                last_page=page_number,
                dpi=180,
                poppler_path=str(poppler_path) if poppler_path else None,
            )
    except Exception as exc:
        raise ParseError("PDF 扫描页渲染失败。请运行 scripts/install-ocr.ps1 安装 Poppler。") from exc

    if not images:
        raise ParseError(f"PDF 第 {page_number} 页无法渲染。")

    return images[0]


def parse_file_to_text(path: Path, ocr_config: OCRConfig | None = None) -> str:
    return parse_file(path, ocr_config).raw_text
